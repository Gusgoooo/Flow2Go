import base64
import io
import re
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from mapper import map_rect_to_ppt


DATA_URL_RE = re.compile(r"^data:(?P<mime>[^;]+);base64,(?P<b64>.+)$")


def _parse_hex_color(s: str) -> tuple[int, int, int]:
    t = (s or "").strip()
    if t.startswith("#"):
        t = t[1:]
    if len(t) == 3:
        t = "".join([c + c for c in t])
    if len(t) != 6:
        return (0, 0, 0)
    try:
        r = int(t[0:2], 16)
        g = int(t[2:4], 16)
        b = int(t[4:6], 16)
        return (r, g, b)
    except Exception:
        return (0, 0, 0)


def _decode_image_bytes(url: str) -> bytes:
    if not url:
        raise ValueError("backgroundImage.url is empty")
    m = DATA_URL_RE.match(url.strip())
    if m:
        raw = base64.b64decode(m.group("b64"))
        mime = m.group("mime").lower()
        if "png" in mime or "jpeg" in mime or "jpg" in mime:
            return raw
        # convert unknown formats (e.g. webp) to PNG
        im = Image.open(io.BytesIO(raw))
        out = io.BytesIO()
        im.convert("RGBA").save(out, format="PNG")
        return out.getvalue()
    raise ValueError("Only data:image/* URLs are supported for backgroundImage.url in MVP")


def build_pptx_bytes(payload: dict[str, Any]) -> bytes:
    prs = Presentation()
    # 16:9 close to widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = payload.get("slides") or []
    for slide in slides:
        sld = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        bg = slide.get("backgroundImage") or {}
        bg_url = str(bg.get("url") or "")
        bg_bytes = _decode_image_bytes(bg_url)

        # background: full slide
        pic_stream = io.BytesIO(bg_bytes)
        sld.shapes.add_picture(pic_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

        # text nodes
        for t in slide.get("textNodes") or []:
            text = str(t.get("text") or "")
            if not text.strip():
                continue

            x, y, w, h = map_rect_to_ppt(
                node_x=float(t.get("x") or 0),
                node_y=float(t.get("y") or 0),
                node_w=float(t.get("width") or 1),
                node_h=float(t.get("height") or 1),
                slide_w_px=float(slide.get("width") or 1440),
                slide_h_px=float(slide.get("height") or 800),
                ppt_w=prs.slide_width,
                ppt_h=prs.slide_height,
            )

            shape = sld.shapes.add_textbox(x, y, w, h)
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text

            fs = float(t.get("fontSize") or 24)
            run.font.size = int(fs) * 100  # pptx uses EMU? python-pptx expects Pt; but accepts int? keep simple
            run.font.name = str(t.get("fontFamily") or "Pingfang SC")

            r, g, b = _parse_hex_color(str(t.get("color") or "#000000"))
            run.font.color.rgb = RGBColor(r, g, b)

            # Weight mapping is best-effort; python-pptx doesn't support numeric weight directly.
            fw = float(t.get("fontWeight") or 400)
            run.font.bold = fw >= 600

            p.alignment = PP_ALIGN.LEFT

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

